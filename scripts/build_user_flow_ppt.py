"""TrendSys 초보자 설명용 PPTX 생성기."""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "TrendSys_사용자용_분류와주제_전체플로우.pptx"

NAVY = RGBColor(18, 28, 48); BLUE = RGBColor(47, 111, 237); CYAN = RGBColor(34, 184, 207)
GREEN = RGBColor(39, 174, 96); ORANGE = RGBColor(242, 153, 74); RED = RGBColor(224, 79, 95)
BG = RGBColor(246, 248, 252); WHITE = RGBColor(255, 255, 255); TEXT = RGBColor(37, 45, 61)
MUTED = RGBColor(102, 112, 133); LINE = RGBColor(218, 224, 235); PURPLE = RGBColor(132, 94, 194)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def box(slide, x, y, w, h, text="", fill=WHITE, color=TEXT, size=18, bold=False,
        radius=True, border=LINE, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    tf = shape.text_frame; tf.clear(); tf.margin_left = tf.margin_right = Inches(.16)
    tf.margin_top = tf.margin_bottom = Inches(.10); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = "Apple SD Gothic Neo"; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color
    return shape

def text(slide, x, y, w, h, value, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=s.text_frame; tf.clear()
    tf.word_wrap=True; tf.margin_left=tf.margin_right=Inches(.02)
    for i, line in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(5)
        r=p.add_run(); r.text=line; r.font.name="Apple SD Gothic Neo"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return s

def base(title, subtitle=None, section=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=BG
    if section: box(s,.45,.35,1.35,.38,section,fill=BLUE,color=WHITE,size=12,bold=True,border=BLUE,align=PP_ALIGN.CENTER)
    text(s,.55,.82,12.2,.55,title,27,NAVY,True)
    if subtitle: text(s,.57,1.38,12,.38,subtitle,13,MUTED)
    text(s,12.35,7.08,.45,.2,str(len(prs.slides)),10,MUTED,False,PP_ALIGN.RIGHT)
    return s

def bullets(slide, x,y,w,h, items, size=17, color=TEXT):
    s=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=s.text_frame; tf.clear(); tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.level=0; p.space_after=Pt(10)
        p.font.name="Apple SD Gothic Neo"; p.font.size=Pt(size); p.font.color.rgb=color
    return s

def arrow(slide,x,y,w=.45,h=.35,color=BLUE):
    return box(slide,x,y,w,h,"→",fill=BG,color=color,size=23,bold=True,border=BG,align=PP_ALIGN.CENTER)

# 1
s=prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
box(s,.7,.65,2.15,.45,"TrendSys 사용자 안내",fill=BLUE,color=WHITE,size=14,bold=True,border=BLUE,align=PP_ALIGN.CENTER)
text(s,.75,1.55,11.8,1.35,"게임 게시판의 수많은 글이\n‘카테고리·주제·알림’이 되는 과정",32,WHITE,True)
text(s,.78,3.2,11.2,.8,"초보자를 위한 전체 플로우 · 분류 근거 · 묶음 기준 · 운영 화면 읽는 법",18,RGBColor(190,205,230))
for i,(a,b,c) in enumerate([("카테고리","글의 큰 성격","서랍"),("주제","지금 함께 떠오른 구체적 이슈","원탁"),("알림","운영팀이 지금 봐야 하는 사건","신호")]):
    box(s,.8+i*4.1,5.15,3.65,1.2,f"{a}\n{b} — {c}",fill=RGBColor(29,43,70),color=WHITE,size=17,bold=True,border=RGBColor(63,82,115),align=PP_ALIGN.CENTER)

# 2
s=base("먼저 결론: 카테고리와 주제는 같은 것이 아닙니다",section="핵심 개념")
box(s,.7,2.05,5.75,3.85,"카테고리 = 큰 분류 서랍\n\n예: 일반 · 콘텐츠 · 운영 · 밸런스 · 과금 · 버그\n\n주된 판단자: Gemma 언어 모델\n벡터는 거의 같은 글을 빠르게 넣는 보조 수단",fill=RGBColor(235,242,255),border=BLUE,size=19,bold=True)
box(s,6.85,2.05,5.75,3.85,"주제 = 지금 함께 이야기되는 구체적 사건\n\n예: ‘신규 캐릭터 스킬 오류’, ‘특정 상품 가격 논란’\n\n주된 판단자: 임베딩 벡터의 가까움\n이름 붙이기·검증·애매한 병합에는 Gemma 사용",fill=RGBColor(235,250,246),border=GREEN,size=19,bold=True)
text(s,.8,6.3,11.8,.45,"한 카테고리 안에 여러 주제가 있을 수 있고, 한 주제에 여러 카테고리 성격이 섞일 수도 있습니다.",16,RED,True,PP_ALIGN.CENTER)

# 3
s=base("한 장으로 보는 전체 흐름", "매 수집 주기와 분석 주기마다 자동으로 반복됩니다.", "전체 흐름")
steps=[("① 수집","원본 DB"),("② 정리","본문 정제"),("③ 벡터화","의미 좌표"),("④ 카테고리","큰 서랍"),("⑤ 주제","구체 이슈"),("⑥ 탐지","열기·급증"),("⑦ 판단","요약·O/X"),("⑧ 표시","UI·이력")]
for i,(a,b) in enumerate(steps):
    x=.45+i*1.59; box(s,x,2.5,1.35,1.35,f"{a}\n{b}",fill=WHITE,border=BLUE if i in (2,3,4) else LINE,size=15,bold=True,align=PP_ALIGN.CENTER)
    if i<7: arrow(s,x+1.34,2.98,.25,.28)
box(s,.8,4.45,11.7,1.25,"두 개의 반복 루프\n수집 루프: 새 글만 가져와 저장·임베딩·카테고리 배정  |  분석 루프: 최근 창을 다시 보고 주제·버스트·알림 갱신",fill=RGBColor(238,241,248),border=LINE,size=17,bold=True,align=PP_ALIGN.CENTER)

# 4
s=base("1단계 — 새 글을 안전하게 가져옵니다", "현재 구현은 dcinside SQLite 원본을 읽기 전용으로 증분 수집하며 mock 모드도 제공합니다.", "수집")
bullets(s,.75,2.0,6.0,3.9,["게임별 원본 DB와 갤러리 필터를 사용합니다.","마지막으로 본 post_no 이후의 글만 가져옵니다.","최초 연결 때는 설정된 시간만큼 소급합니다(기본 72시간).","글마다 game_id를 붙여 서로 다른 게임 데이터가 섞이지 않게 합니다.","수집 커서는 DB에 저장되어 서버 재시작 후에도 이어집니다."],17)
box(s,7.15,2.25,5.3,.8,"원본 DB → 새 번호만 조회",fill=RGBColor(235,242,255),border=BLUE,size=20,bold=True,align=PP_ALIGN.CENTER)
arrow(s,9.55,3.25,.45,.4)
box(s,7.15,3.85,5.3,1.1,"TrendSys DB\nposts(game_id, id, title, body, created_at…)",fill=WHITE,border=GREEN,size=18,bold=True,align=PP_ALIGN.CENTER)

# 5
s=base("2단계 — 분석하기 좋은 문장으로 정리합니다", "원문을 버리는 것이 아니라, 분석용 text를 별도로 만듭니다.", "전처리")
box(s,.75,2.05,5.55,3.6,"입력 예시\n\n‘ㅋㅋㅋ &amp; 이번 패치 ㄹㅇ...’\n긴 본문 / HTML / 반복 공백 / 게시판 표현",fill=WHITE,border=ORANGE,size=19,bold=True)
arrow(s,6.38,3.45,.55,.45)
box(s,7.05,2.05,5.55,3.6,"정리 결과\n\nHTML 해제 · 공백 정리 · 은어 사전 적용\n너무 짧은 글 제외 · 최대 길이 제한\n제목 + 본문을 분석용 텍스트로 구성",fill=RGBColor(255,248,235),border=ORANGE,size=19,bold=True)
text(s,.8,6.15,11.7,.55,"이미지가 있으면 캡션을 덧붙일 수 있지만, 현재 기본 캡셔닝은 준비 상태에 따라 stub일 수 있습니다.",15,MUTED,False,PP_ALIGN.CENTER)

# 6
s=base("3단계 — 글의 의미를 숫자 좌표로 바꿉니다", "임베딩은 ‘뜻이 비슷한 문장은 지도에서 가까운 점’이 되도록 만드는 기술입니다.", "임베딩")
box(s,.7,2.0,4.05,1.05,"‘신캐 스킬이 작동 안 함’",fill=WHITE,border=RED,size=18,bold=True,align=PP_ALIGN.CENTER)
box(s,.7,3.35,4.05,1.05,"‘새 캐릭터 기술 버그’",fill=WHITE,border=RED,size=18,bold=True,align=PP_ALIGN.CENTER)
box(s,.7,4.7,4.05,1.05,"‘오늘 점심 뭐 먹지’",fill=WHITE,border=MUTED,size=18,bold=True,align=PP_ALIGN.CENTER)
text(s,5.0,2.35,1.5,2.7,"→\n\n→\n\n→",23,BLUE,True,PP_ALIGN.CENTER)
box(s,6.35,1.95,5.95,4.0,"의미 지도(1024차원)\n\n● A  ● B   ← 서로 가까움\n\n                         ● C ← 멀리 떨어짐\n\n실제 UI에 좌표를 그리는 것은 아니며, 코사인 유사도로 방향의 가까움을 계산합니다.",fill=RGBColor(239,247,255),border=CYAN,size=18,bold=True)
text(s,.8,6.3,11.7,.4,"중요: 임베딩은 새 글마다 한 번만 계산해 DB에 저장합니다. 모델을 바꾸면 전량 재계산과 임계값 재조정이 필요합니다.",14,RED,True,PP_ALIGN.CENTER)

# 7
s=base("코사인 유사도는 무엇인가요?", "두 의미 벡터가 같은 방향을 보는 정도: 1에 가까울수록 비슷합니다.", "기초 개념")
for i,(v,label,c) in enumerate([("0.95","거의 같은 글·도배",GREEN),("0.85","상당히 유사",BLUE),("0.70","주제 연결 후보",ORANGE),("0.30","관련성 낮음",MUTED)]):
    box(s,.9+i*3.05,2.2,2.6,1.4,f"{v}\n{label}",fill=WHITE,border=c,size=19,bold=True,align=PP_ALIGN.CENTER)
box(s,.85,4.35,11.65,1.35,"주의: 숫자의 의미는 임베딩 모델과 실제 데이터에 따라 달라집니다.\n이 프로그램의 값은 KURE-v1과 실제 24시간 게시글 표본을 기준으로 보정한 운영 설정입니다.",fill=RGBColor(255,247,235),border=ORANGE,size=18,bold=True,align=PP_ALIGN.CENTER)

# 8
s=base("4단계 — 카테고리는 ‘글의 큰 성격’을 분류합니다", "카테고리는 계속 유지되는 자산이며, 게임별로 분리됩니다.", "카테고리")
for i,(name,desc) in enumerate([("일반","잡담·질문"),("콘텐츠","캐릭터·공략"),("운영","공지·제재"),("밸런스","성능·너프"),("과금","상품·가격"),("버그","오류·장애")]):
    x=.75+(i%3)*4.15; y=2.0+(i//3)*1.45
    box(s,x,y,3.75,1.05,f"{name}  |  {desc}",fill=WHITE,border=BLUE,size=18,bold=True,align=PP_ALIGN.CENTER)
box(s,.9,5.25,11.5,1.05,"카테고리 ID·이름·중심벡터·누적 글 수·마지막 알림 시각이 DB에 남습니다.\n분석을 다시 실행해도 번호가 매번 바뀌는 임시 클러스터가 아닙니다.",fill=RGBColor(235,242,255),border=BLUE,size=16,bold=True,align=PP_ALIGN.CENTER)

# 9
s=base("카테고리 배정의 실제 근거: 2단계 하이브리드", "코드상 assign_posts()가 새 글마다 아래 순서를 따릅니다.", "카테고리")
box(s,.7,2.0,3.5,3.95,"1. 벡터 빠른 경로\n\n기존 카테고리 중심과 비교\nCATEGORY_ASSIGN_SIM 이상이면 즉시 배정\n\n현재 설계상 이 경로는 ‘거의 같은 글’ 위주\nLLM 호출을 줄이는 지름길",fill=RGBColor(235,250,246),border=GREEN,size=18,bold=True)
arrow(s,4.35,3.65,.55,.45)
box(s,5.0,2.0,3.5,3.95,"2. Gemma 의미 판단\n\n유사한 기존 카테고리 후보 K개와 글을 함께 제시\n\n기존 이름을 고르거나\n필요하면 새 카테고리 이름 생성\n\n언어 의미가 주 근거",fill=RGBColor(241,237,255),border=PURPLE,size=18,bold=True)
arrow(s,8.65,3.65,.55,.45)
box(s,9.3,2.0,3.3,3.95,"3. 영속 저장\n\npost.category_id 기록\n카테고리 중심 갱신\npost_count 증가\n\n중심 갱신은 상한에서 멈춰 ‘잡담 블랙홀’로 표류하는 것을 방지",fill=RGBColor(235,242,255),border=BLUE,size=18,bold=True)

# 10
s=base("왜 카테고리를 벡터만으로 정하지 않나요?", "실측 결과, 큰 분류의 의미 경계는 벡터 거리만으로 잘 갈리지 않았습니다.", "설계 근거")
box(s,.8,2.0,5.7,3.75,"실측\n\n글 ↔ 자기 카테고리 중심: 약 0.661\n글 ↔ 다른 카테고리 중심: 약 0.653\n\n차이가 너무 작아 임계값을 어디에 두어도\n통과율과 오배정률이 함께 움직였습니다.",fill=WHITE,border=RED,size=20,bold=True,align=PP_ALIGN.CENTER)
box(s,6.85,2.0,5.7,3.75,"결론\n\n‘버그인가, 밸런스 불만인가?’는\n단어의 가까움보다 문맥 이해가 중요합니다.\n\n따라서 카테고리의 주 배정자는 Gemma,\n벡터는 근사 중복용 빠른 경로입니다.",fill=RGBColor(241,237,255),border=PURPLE,size=20,bold=True,align=PP_ALIGN.CENTER)

# 11
s=base("5단계 — 주제는 ‘지금 같이 떠오른 구체 이슈’를 묶습니다", "최근 분석 창(기본 24시간)의 글을 의미 벡터로 묶습니다.", "주제")
box(s,.7,2.0,3.55,3.9,"① 가까운 점 연결\n\nTOPIC_LINK_SIM 기본 0.70\n탐욕적 그룹화로 가까운 글을 같은 묶음에 연결\n\n최소 TOPIC_MIN_SIZE 기본 4건 이상이어야 주제로 인정",fill=RGBColor(235,250,246),border=GREEN,size=18,bold=True)
arrow(s,4.35,3.65,.5,.4)
box(s,4.95,2.0,3.55,3.9,"② 기존 주제 추적\n\n묶음 중심과 저장된 주제 중심 비교\nTOPIC_MATCH_SIM 기본 0.85 이상이면 기존 topic_id 유지\n\n아니면 신규 주제로 생성",fill=RGBColor(235,242,255),border=BLUE,size=18,bold=True)
arrow(s,8.6,3.65,.5,.4)
box(s,9.2,2.0,3.4,3.9,"③ 이름과 소속 정제\n\nGemma가 대표 글을 보고 이름 부여\n멤버 글 O/X 검증\n표류 시 재명명\n갈라진 주제는 자동 또는 LLM 병합",fill=RGBColor(241,237,255),border=PURPLE,size=18,bold=True)

# 12
s=base("주제 묶음이 틀어지는 것을 어떻게 막나요?", "실시간 게시판은 잡담·도배·대화 표류가 있어 정화 장치가 필요합니다.", "품질 관리")
items=[("멤버 검증","주제명과 각 글이 실제로 맞는지 Gemma가 O/X 판정. X는 주제 상세에서도 제외"),
       ("드리프트 재명명","현재 중심이 이름을 붙였던 중심과 0.85 미만으로 멀어지면 이름을 다시 생성"),
       ("자동 병합","두 주제 중심 유사도가 0.90 이상이면 같은 주제로 자동 병합"),
       ("애매한 병합","0.75 이상 0.90 미만은 Gemma가 같은 사건인지 판정하고 결과를 캐시"),
       ("생성 시 편입","새 묶음이 기존 주제와 애매하게 가깝다면 이름 생성 단계에서 기존 사건 편입 여부 판단"),
       ("HDBSCAN 안전망","별도 유지보수 도구로 카테고리 중복·누락 후보를 점검")]
for i,(a,b) in enumerate(items):
    x=.7+(i%2)*6.15; y=1.95+(i//2)*1.55
    box(s,x,y,5.75,1.2,f"{a}\n{b}",fill=WHITE,border=[GREEN,BLUE,PURPLE,ORANGE,BLUE,RED][i],size=15,bold=True)

# 13
s=base("6단계 — ‘많다’가 아니라 ‘갑자기 뜨겁다’를 찾습니다", "각 카테고리와 주제에 시간 감쇠 열기(heat)를 계산합니다.", "버스트")
box(s,.75,2.0,3.55,3.8,"시간 감쇠\n\n방금 올라온 글 = 큰 가중치\n오래된 글 = 작은 가중치\n\n기본 반감기 30분\n30분 전 글은 현재 글의 절반 정도 영향",fill=RGBColor(255,248,235),border=ORANGE,size=19,bold=True,align=PP_ALIGN.CENTER)
box(s,4.85,2.0,3.55,3.8,"급증 조건\n\n최근 창 기본 60분\n평소 대비 기본 3배\n최소 최근 글 기본 5건\n\n비율 + 절대량을 함께 봐 작은 표본 오탐 방지",fill=RGBColor(255,239,241),border=RED,size=19,bold=True,align=PP_ALIGN.CENTER)
box(s,8.95,2.0,3.55,3.8,"대상 구분\n\n카테고리 heat: 화면의 큰 흐름\n주제 heat: 실제 알림 판단 기준\n\n버스트가 아닌 주제는 비싼 요약·판단을 실행하지 않음",fill=RGBColor(235,242,255),border=BLUE,size=19,bold=True,align=PP_ALIGN.CENTER)

# 14
s=base("7단계 — 급증한 주제만 요약하고 알림 O/X를 판단합니다", "알림 직전에는 멤버를 다시 검증해 섞인 글이 판단을 오염시키지 않게 합니다.", "알림")
flow=[("버스트","급증 확인",RED),("재검증","이물질 제거",ORANGE),("요약","대표 글 최대 15건",CYAN),("라벨","adapters4 다수결",PURPLE),("판단","adapters3 O/X",BLUE),("발송","쿨다운 확인",GREEN)]
for i,(a,b,c) in enumerate(flow):
    x=.5+i*2.1; box(s,x,2.3,1.75,1.4,f"{a}\n{b}",fill=WHITE,border=c,size=15,bold=True,align=PP_ALIGN.CENTER)
    if i<5: arrow(s,x+1.77,2.78,.3,.3)
box(s,.75,4.45,11.85,1.35,"판단 입력 형식:  주제 라벨 ___ 요약 제목 ___ 요약 본문\nO = 운영팀에 즉시 알릴 가치가 있음  |  X = 지금은 알리지 않음\nO와 X 모두 이력에 남기고, 담당자 피드백을 다음 파인튜닝 데이터로 내보낼 수 있습니다.",fill=RGBColor(241,237,255),border=PURPLE,size=17,bold=True,align=PP_ALIGN.CENTER)

# 15
s=base("알림 이후에는 ‘소강’도 감시합니다", "이슈가 끝났다는 신호는 LLM 없이 수치로 판단합니다.", "종결")
box(s,.85,2.1,3.45,3.6,"알림 시점\n\n주제의 최고 heat를 기록\nlast_alerted_at 저장\n같은 주제 재알림은 기본 60분 쿨다운",fill=RGBColor(255,239,241),border=RED,size=18,bold=True,align=PP_ALIGN.CENTER)
arrow(s,4.45,3.55,.55,.4)
box(s,5.1,2.1,3.45,3.6,"소강 조건\n\n현재 heat ≤ 정점의 25%\nAND 최근 글 < 최소 글 수\n조건이 기본 10사이클 연속 유지",fill=RGBColor(255,248,235),border=ORANGE,size=18,bold=True,align=PP_ALIGN.CENTER)
arrow(s,8.7,3.55,.55,.4)
box(s,9.35,2.1,3.1,3.6,"종료 신호\n\n‘정점 → 현재 heat’와 최근 글 수를 요약해 소강 알림 저장·발송\n플래핑 방지",fill=RGBColor(235,250,246),border=GREEN,size=18,bold=True,align=PP_ALIGN.CENTER)

# 16
s=base("사용자 화면에서는 이렇게 읽으면 됩니다", "Vue 화면은 10~15초 간격으로 최신 서버 상태를 다시 읽습니다.", "화면 안내")
for i,(a,b,c) in enumerate([("실시간 관제","슬롯 상태·운영 다이얼·카테고리/주제 열기·최근 글",BLUE),("주제","구체 이슈 목록. 알림 → 버스트 → heat 순",GREEN),("카테고리","큰 분류의 누적 흐름과 스파크라인",CYAN),("알림","현재 O 판정과 O/X 이력·담당자 피드백",RED),("타임머신","과거 스냅샷 재생·구간 알림/상위 주제",PURPLE),("점검 도구","HDBSCAN·코사인 실험·Gemma 분석·판단 테스트",ORANGE)]):
    x=.65+(i%2)*6.25; y=1.9+(i//2)*1.55
    box(s,x,y,5.85,1.18,f"{a}\n{b}",fill=WHITE,border=c,size=16,bold=True)

# 17
s=base("결과를 해석할 때 꼭 알아야 할 한계", "이 시스템은 ‘자동 정답기’가 아니라 운영자가 빠르게 볼 후보를 만드는 관제 도구입니다.", "주의사항")
bullets(s,.8,1.95,11.8,4.8,["유사도 임계값은 보편적 진리가 아닙니다. 임베딩 모델·게임·게시판 문화가 바뀌면 재보정해야 합니다.","Gemma 분류·이름·검증·알림 판단은 오분류할 수 있습니다. 알림 화면의 피드백이 개선 데이터가 됩니다.","카테고리는 의미 중심, 주제는 벡터 중심이라 서로 1:1 관계가 아닙니다.","글이 4건 미만이면 의미가 있어도 주제로 표시되지 않을 수 있습니다.","새로운 이슈는 기존 주제와 0.85 미만이면 새 ID로 생길 수 있고, 이후 병합될 수 있습니다.","실제 외부 알림 채널 연결 지점은 현재 send_alert()이며 기본 구현은 로그 출력입니다.","이미지 캡셔닝·요약은 설정된 backend 준비 상태에 따라 stub으로 동작할 수 있습니다."],17)

# 18
s=base("운영자가 조정하는 핵심 다이얼", "한 값을 바꾸면 정확도·속도·비용 사이의 균형이 달라집니다.", "설정")
rows=[("CATEGORY_ASSIGN_SIM","카테고리 벡터 빠른 배정","낮추면 오배정↑ / 높이면 LLM 호출↑"),
      ("TOPIC_LINK_SIM = 0.70","글끼리 주제 연결","낮추면 큰 잡탕 / 높이면 잘게 분열"),
      ("TOPIC_MATCH_SIM = 0.85","기존 주제 ID 유지","낮추면 흡수 / 높이면 새 주제 증가"),
      ("TOPIC_AUTO_MERGE_SIM = 0.90","자동 병합","낮추면 블랙홀 병합 위험"),
      ("HALF_LIFE_MIN = 30","오래된 글 영향 감소","짧으면 민감 / 길면 오래 뜨거움"),
      ("BURST_RATIO = 3 · MIN_RECENT = 5","급증 판정","낮추면 민감·오탐 / 높이면 보수적"),
      ("COOLDOWN_MIN = 60","같은 주제 재알림 제한","짧으면 반복 / 길면 재점화 놓침")]
for i,(a,b,c) in enumerate(rows):
    y=1.85+i*.68
    box(s,.65,y,3.65,.52,a,fill=RGBColor(235,242,255),border=BLUE,size=13,bold=True)
    box(s,4.4,y,3.55,.52,b,fill=WHITE,border=LINE,size=13,bold=True)
    box(s,8.05,y,4.65,.52,c,fill=WHITE,border=LINE,size=13)

# 19
s=base("코드 기준: 설명의 근거가 되는 주요 위치", "PPT 내용은 현재 작업 트리의 실제 구현을 기준으로 작성했습니다.", "참고")
refs=[("app/main.py","주기 수집·분석 루프, Gemma 큐 워커"),("app/services/pipeline.py","전체 연결, heat·버스트·요약·판정·소강"),("app/services/categorize.py","카테고리 벡터 빠른 경로 + Gemma 분류"),("app/services/topics.py","주제 묶기·영속 추적·검증·재명명·병합"),("app/services/detection.py","시간 감쇠와 버스트 공식"),("app/services/gemma_analyze.py","글 단위 대분류·주제 라벨·감성"),("app/db.py","게시글·카테고리·주제·점수·스냅샷·알림 저장"),("app/routers/analysis.py","사용자 화면이 호출하는 API"),("frontend/src","Vue MVVM 화면·폴링·게임 선택")]
for i,(a,b) in enumerate(refs):
    x=.7+(i%2)*6.15; y=1.85+(i//2)*.91
    box(s,x,y,2.55,.7,a,fill=RGBColor(238,241,248),border=LINE,size=13,bold=True)
    box(s,x+2.65,y,3.1,.7,b,fill=WHITE,border=LINE,size=13)

# 20
s=base("마지막으로 기억할 세 문장",section="요약")
box(s,.9,1.95,11.55,1.15,"1. 카테고리는 글의 ‘큰 성격’이며, 문맥을 이해하는 Gemma가 주로 결정합니다.",fill=RGBColor(235,242,255),border=BLUE,size=21,bold=True,align=PP_ALIGN.CENTER)
box(s,.9,3.35,11.55,1.15,"2. 주제는 ‘지금 같이 떠오른 구체 이슈’이며, 임베딩 벡터의 가까움으로 묶고 Gemma가 이름·검증·병합을 돕습니다.",fill=RGBColor(235,250,246),border=GREEN,size=20,bold=True,align=PP_ALIGN.CENTER)
box(s,.9,4.75,11.55,1.15,"3. 알림은 급증한 주제만 요약·판단하며, 피드백과 이력으로 운영 품질을 계속 개선합니다.",fill=RGBColor(241,237,255),border=PURPLE,size=21,bold=True,align=PP_ALIGN.CENTER)

OUT.parent.mkdir(parents=True, exist_ok=True); prs.save(OUT); print(OUT)
