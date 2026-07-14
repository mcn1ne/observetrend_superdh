"""TrendSys 전체 작동 방식 PPTX 생성기 — 기획자/사용자용 (2026-07-13 코드 기준).

실행:  uv run --with python-pptx python scripts/build_system_overview_ppt.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "docs" / "TrendSys_전체작동방식.pptx"

NAVY = RGBColor(18, 28, 48); BLUE = RGBColor(47, 111, 237); CYAN = RGBColor(34, 184, 207)
GREEN = RGBColor(39, 174, 96); ORANGE = RGBColor(242, 153, 74); RED = RGBColor(224, 79, 95)
BG = RGBColor(246, 248, 252); WHITE = RGBColor(255, 255, 255); TEXT = RGBColor(37, 45, 61)
MUTED = RGBColor(102, 112, 133); LINE = RGBColor(218, 224, 235); PURPLE = RGBColor(132, 94, 194)
L_BLUE = RGBColor(235, 242, 255); L_GREEN = RGBColor(235, 250, 246); L_PURPLE = RGBColor(241, 237, 255)
L_ORANGE = RGBColor(255, 248, 235); L_RED = RGBColor(255, 239, 241); L_GRAY = RGBColor(238, 241, 248)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)


def box(slide, x, y, w, h, text="", fill=WHITE, color=TEXT, size=18, bold=False,
        radius=True, border=LINE, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    tf = shape.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.14)
    tf.margin_top = tf.margin_bottom = Inches(.08); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line; r.font.name = "Apple SD Gothic Neo"
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return shape


def text(slide, x, y, w, h, value, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(.02)
    for i, line in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(5)
        r = p.add_run(); r.text = line; r.font.name = "Apple SD Gothic Neo"
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return s


def base(title, subtitle=None, section=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = BG
    if section:
        box(s, .45, .35, 1.55, .38, section, fill=BLUE, color=WHITE, size=12,
            bold=True, border=BLUE, align=PP_ALIGN.CENTER)
    text(s, .55, .82, 12.2, .55, title, 26, NAVY, True)
    if subtitle:
        text(s, .57, 1.38, 12.2, .38, subtitle, 13, MUTED)
    text(s, 12.35, 7.08, .45, .2, str(len(prs.slides)), 10, MUTED, False, PP_ALIGN.RIGHT)
    return s


def arrow(slide, x, y, w=.4, h=.35, label="→", color=BLUE):
    return box(slide, x, y, w, h, label, fill=BG, color=color, size=21, bold=True,
               border=BG, align=PP_ALIGN.CENTER)


def down_arrow(slide, x, y, color=BLUE):
    return box(slide, x, y, .4, .32, "↓", fill=BG, color=color, size=20, bold=True,
               border=BG, align=PP_ALIGN.CENTER)


# ── 1. 표지 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
box(s, .7, .65, 2.6, .45, "TrendSys 작동 방식 안내", fill=BLUE, color=WHITE, size=14,
    bold=True, border=BLUE, align=PP_ALIGN.CENTER)
text(s, .75, 1.5, 12, 1.5, "게시판의 수천 개 글이\n'알림 한 줄'이 되기까지", 34, WHITE, True)
text(s, .78, 3.25, 11.4, .5, "전체 파이프라인을 순서대로 · 흐름도 중심 · 개발 지식 불필요  (2026-07-13 코드 기준)",
     17, RGBColor(190, 205, 230))
for i, (a, b) in enumerate([("읽는다", "매분 수집해 AI가 다 읽음"),
                            ("묶는다", "같은 이야기끼리 자동으로"),
                            ("알린다", "갑자기 커진 사건만 골라서")]):
    box(s, .8 + i * 4.1, 4.9, 3.7, 1.5, f"{a}\n{b}", fill=RGBColor(29, 43, 70), color=WHITE,
        size=17, bold=True, border=RGBColor(63, 82, 115), align=PP_ALIGN.CENTER)

# ── 2. 무엇을 하는 시스템인가 ───────────────────────────────────────
s = base("이 시스템이 풀려는 문제", "게시판 여론은 빠르고, 사람은 다 읽을 수 없습니다.", "왜 필요한가")
box(s, .7, 2.0, 5.8, 3.9,
    "문제\n\n· 갤러리에 하루 800~2,400개의 글\n· 반어·밈·욕설이 섞여 그대로 읽기 어려움\n"
    "· 진짜 사건(버그·결제 장애·민심 폭발)은\n  잡담 사이에 묻혀서 뒤늦게 발견됨",
    fill=L_RED, border=RED, size=18, bold=True)
box(s, 6.85, 2.0, 5.8, 3.9,
    "해법 — 24시간 자동 관제\n\n· 매분 새 글을 수집해 AI가 전부 읽음\n"
    "· 글을 '종류(카테고리)'와 '사건(주제)'으로 정리\n"
    "· 갑자기 커지는 사건만 요약해서 알림\n· 모든 AI는 이 컴퓨터 안에서 동작 (외부 유출 없음)",
    fill=L_GREEN, border=GREEN, size=18, bold=True)
text(s, .8, 6.15, 11.8, .5, "핵심 철학: 값싼 계산(벡터·규칙)으로 후보를 좁히고, 비싼 AI 판단은 좁힌 곳에만 씁니다.",
     16, NAVY, True, PP_ALIGN.CENTER)

# ── 3. 전체 흐름 한 장 ──────────────────────────────────────────────
s = base("전체 흐름 한 장", "글 하나가 들어와서 알림이 되기까지 — 아래 순서대로 이 자료가 설명합니다.", "큰 그림")
steps = [("① 수집", "매분 새 글", BLUE), ("② 정리", "은어·중복 정돈", BLUE),
         ("③ 좌표화", "뜻을 숫자로", CYAN), ("④ 글 읽기", "종류·감성 분류", PURPLE),
         ("⑤ 주제 묶기", "같은 사건끼리", GREEN), ("⑥ 급증 감지", "갑자기 뜨거운가", ORANGE),
         ("⑦ 판단", "알릴 가치 O/X", RED), ("⑧ 알림·화면", "관제 대시보드", NAVY)]
for i, (a, b, c) in enumerate(steps):
    x = .42 + i * 1.62
    box(s, x, 2.35, 1.42, 1.5, f"{a}\n{b}", fill=WHITE, border=c, size=14, bold=True,
        align=PP_ALIGN.CENTER)
    if i < 7:
        arrow(s, x + 1.4, 2.92, .25, .3)
box(s, .7, 4.45, 6.0, 1.9,
    "쉬지 않는 일꾼 4명 (백그라운드 루프)\n"
    "· 수집 루프: 매분 새 글 가져와 저장·좌표화\n"
    "· 분석 루프: 매분 주제·급증·알림 다시 계산\n"
    "· 글 읽기 워커: 밀린 글을 순서대로 정독\n"
    "· 이미지 워커: 첨부 사진을 읽어 글에 반영",
    fill=L_GRAY, border=LINE, size=15, bold=True)
box(s, 7.0, 4.45, 5.6, 1.9,
    "글이 몰려도 안 잃어버립니다\n\n읽기·이미지 작업은 '대기열(큐)'에 쌓였다가\n"
    "한가할 때 자동으로 따라잡습니다.\n알림 감지는 그동안에도 계속 돕니다.",
    fill=L_BLUE, border=BLUE, size=15, bold=True)

# ── 4. 등장 AI 소개 ────────────────────────────────────────────────
s = base("등장하는 AI — 전부 이 컴퓨터 안에서 돕니다", "역할별로 다른 모델(또는 같은 모델의 다른 훈련본)을 씁니다. 외부 API·과금 없음.", "등장인물")
crew = [("좌표 변환기", "KURE-v1 임베딩", "글의 뜻을 1,024개 숫자 좌표로 변환 — '비슷한 글=가까운 점'", CYAN),
        ("글 읽기 담당", "Gemma 4 + 자체학습 ➍", "글마다 주제문구·종류(6종)·감성(7종)을 뽑음 (글당 약 1초)", PURPLE),
        ("사회자", "Gemma 4 기본형", "주제 이름 짓기 · 멤버 검증 · 병합 판정 · 요약", BLUE),
        ("알림 판단관", "Gemma 4 + 자체학습 ➌", "요약을 보고 '운영팀에 알릴 사건인가' O/X 한 글자로 판단", RED),
        ("눈(이미지)", "Gemma 4 멀티모달", "첨부 사진을 보고 글과 관련 있으면 한 문장으로 묘사", GREEN)]
for i, (a, b, c, col) in enumerate(crew):
    y = 1.9 + i * 1.02
    box(s, .7, y, 2.3, .84, a, fill=col, color=WHITE, size=16, bold=True, align=PP_ALIGN.CENTER)
    box(s, 3.1, y, 2.9, .84, b, fill=WHITE, border=col, size=14, bold=True, align=PP_ALIGN.CENTER)
    box(s, 6.1, y, 6.5, .84, c, fill=WHITE, border=LINE, size=14)
text(s, .8, 7.0, 11.8, .4, "➍➌ = 우리 게시판 데이터로 따로 학습시킨 전용 모델 (GreatestStep 프로젝트)", 12, MUTED)

# ── 5. ① 수집 ──────────────────────────────────────────────────────
s = base("① 수집 — 매분, 새 글만 가져옵니다", "크롤러가 모아둔 원본 DB에서 '지난번 이후 글'만 증분으로 읽습니다.", "1단계 · 수집")
box(s, .7, 2.2, 3.3, 2.6, "크롤러 원본 DB\n\n디시인사이드 갤러리 글\n(제목·본문·시각·첨부)", fill=WHITE,
    border=LINE, size=16, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 4.15, 3.25, .55, .45)
box(s, 4.85, 2.2, 3.6, 2.6, "매분 1회\n\n'마지막으로 본 글 번호'\n이후의 새 글만 읽음\n(재시작해도 이어서)",
    fill=L_BLUE, border=BLUE, size=16, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 8.6, 3.25, .55, .45)
box(s, 9.3, 2.2, 3.3, 2.6, "TrendSys 저장소\n\n게임별로 분리 저장\n(세븐나이츠 리버스,\n나혼자만레벨업)", fill=L_GREEN,
    border=GREEN, size=16, bold=True, align=PP_ALIGN.CENTER)
box(s, .8, 5.25, 11.8, 1.1,
    "여러 게임을 한 화면에서 — 글마다 game_id 를 붙여 게임별 데이터가 절대 섞이지 않게 하고,\n"
    "화면 좌측의 게임 선택기로 전환합니다. 원본은 읽기 전용으로만 열어 크롤러와 충돌하지 않습니다.",
    fill=L_GRAY, border=LINE, size=15, bold=True, align=PP_ALIGN.CENTER)

# ── 6. ② 정리 ──────────────────────────────────────────────────────
s = base("② 정리 — 분석하기 좋은 문장으로 다듬습니다", "원문은 그대로 두고, 분석용 텍스트를 따로 만듭니다.", "2단계 · 전처리")
box(s, .7, 2.1, 5.55, 3.3,
    "들어온 원문\n\n\"ㄹㅇ 이번 패치 &amp; 튕김 실화냐...\n[첨부] https://dcimg8.../viewimage...\"\n\nHTML 찌꺼기 · 은어 · 이미지 링크 · 앱 푸터",
    fill=WHITE, border=ORANGE, size=16, bold=True)
arrow(s, 6.35, 3.5, .55, .45)
box(s, 7.05, 2.1, 5.55, 3.3,
    "정리 결과\n\n\"진짜 이번 패치 강제 종료 실화냐\"\n\n· 은어 사전 적용 (ㄹㅇ→진짜, 튕김→강제 종료)\n· 너무 짧은 글·완전 중복 제외\n· 첨부 이미지 주소는 따로 보관 → ⑤에서 사용",
    fill=L_ORANGE, border=ORANGE, size=16, bold=True)
text(s, .8, 5.85, 11.8, .5, "은어 사전은 운영하며 계속 키웁니다 — \"같은 얘기인데 안 묶인 글\"이 발견될 때마다 원인 은어를 추가.",
     15, MUTED, False, PP_ALIGN.CENTER)

# ── 7. ③ 좌표화 ────────────────────────────────────────────────────
s = base("③ 좌표화 — 글의 '뜻'을 지도 위의 점으로", "컴퓨터는 글자를 비교하지 못하므로, 뜻이 비슷한 글이 가까운 점이 되도록 숫자 좌표로 바꿉니다.", "3단계 · 임베딩")
box(s, .7, 2.0, 4.0, 1.0, "\"신캐 스킬이 작동 안 함\"", fill=WHITE, border=RED, size=16, bold=True, align=PP_ALIGN.CENTER)
box(s, .7, 3.25, 4.0, 1.0, "\"새 캐릭터 기술 버그\"", fill=WHITE, border=RED, size=16, bold=True, align=PP_ALIGN.CENTER)
box(s, .7, 4.5, 4.0, 1.0, "\"오늘 점심 뭐 먹지\"", fill=WHITE, border=MUTED, size=16, bold=True, align=PP_ALIGN.CENTER)
text(s, 4.9, 2.3, 1.2, 3.0, "→\n\n→\n\n→", 22, CYAN, True, PP_ALIGN.CENTER)
box(s, 6.15, 1.95, 6.4, 3.75,
    "의미 지도 (1,024차원)\n\n      ● 스킬 작동 안 함\n      ● 기술 버그        ← 서로 붙어 있음\n\n\n"
    "                                    ● 점심 얘기  ← 멀리 떨어짐\n\n"
    "'가까움'은 두 점의 방향이 얼마나 같은지(코사인 유사도, 0~1)로 잽니다.",
    fill=RGBColor(239, 247, 255), border=CYAN, size=16, bold=True)
text(s, .8, 6.1, 11.8, .5, "글마다 딱 한 번 계산해 저장합니다. 이후의 모든 '비슷하다' 판단이 이 좌표 위에서 이뤄집니다.",
     15, NAVY, True, PP_ALIGN.CENTER)

# ── 8. ④ 글 읽기 → 카테고리 ────────────────────────────────────────
s = base("④ 글 읽기 — AI가 모든 글을 한 번씩 정독합니다", "우리 게시판 데이터로 따로 학습시킨 모델이 글마다 3가지를 뽑고, 그 결과가 카테고리(서랍)를 결정합니다.", "4단계 · 분류")
box(s, .6, 2.15, 2.6, 2.9, "게시글 1건\n\n제목 + 본문", fill=WHITE, border=LINE, size=17, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 3.3, 3.35, .45, .45)
box(s, 3.85, 2.15, 3.1, 2.9, "글 읽기 AI\n(자체학습 ➍)\n\n글당 약 1초\n반어·밈도 학습으로 이해",
    fill=L_PURPLE, border=PURPLE, size=16, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 7.05, 3.35, .45, .45)
for i, (a, b, c) in enumerate([("주제문구", "\"[밸런스] 루디 반격 성능 불만\"", GREEN),
                               ("종류(대분류)", "일반·콘텐츠·운영·밸런스·과금·버그 중 1", BLUE),
                               ("감성", "중립·불만·긍정 등 7가지 중 1", ORANGE)]):
    box(s, 7.6, 2.05 + i * 1.08, 5.05, .92, f"{a} — {b}", fill=WHITE, border=c, size=14, bold=True)
box(s, .8, 5.5, 11.8, 1.2,
    "종류(대분류) → 카테고리 서랍 배정은 고정 규칙표로 즉시 끝납니다 (추가 AI 호출 없음).\n"
    "'일반' 판정만 주제문구에 육성·장비·덱 표현이 있으면 [캐릭터·장비·덱], 없으면 [일반·잡담]으로 나눕니다.",
    fill=L_BLUE, border=BLUE, size=15, bold=True, align=PP_ALIGN.CENTER)

# ── 9. 카테고리 = 고정된 7개 서랍 ──────────────────────────────────
s = base("카테고리 — 게임당 7개로 고정된 '분류 서랍'", "서랍은 지형(전체 분포)을 보여주고, 개별 사건 추적은 다음 장의 '주제'가 담당합니다.", "카테고리")
cats = [("일반·잡담", MUTED), ("캐릭터·장비·덱", GREEN), ("콘텐츠·공략", CYAN), ("밸런스", BLUE),
        ("운영·이벤트", PURPLE), ("과금", ORANGE), ("버그·오류", RED)]
for i, (a, c) in enumerate(cats):
    x = .55 + (i % 4) * 3.13; y = 2.1 + (i // 4) * 1.15
    box(s, x, y, 2.9, .95, a, fill=WHITE, border=c, size=17, bold=True, align=PP_ALIGN.CENTER)
box(s, .7, 4.7, 5.9, 1.85,
    "왜 7개로 고정했나요?\n\n예전에는 AI가 서랍 이름을 자유롭게 만들었더니\n36개까지 불어나고 '유머·잡담' 서랍이\n전체 글의 58%를 빨아들이는 사고가 있었습니다.",
    fill=L_RED, border=RED, size=14, bold=True)
box(s, 6.75, 4.7, 5.9, 1.85,
    "서랍이 보여주는 것\n\n서랍별 글 수·시간별 열기 그래프 → 게시판의 큰 지형\n"
    "\"버그 서랍이 갑자기 붐빈다\" 같은 흐름을 한눈에.\n세부 사건 이름은 주제 계층이 담당합니다.",
    fill=L_BLUE, border=BLUE, size=14, bold=True)

# ── 10. ⑤ 이미지 ───────────────────────────────────────────────────
s = base("④-보조. 이미지 — 옆길에서 읽고, 관련 있을 때만 반영", "글 처리를 막지 않도록 이미지는 별도 일꾼이 나중에 읽습니다 (글당 최대 3장, 약 2초).", "이미지")
box(s, .6, 2.0, 3.2, 1.0, "글 (첨부 3장)", fill=WHITE, border=LINE, size=16, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 3.9, 2.3, .45, .4)
box(s, 4.45, 2.0, 4.1, 1.0, "글은 기다리지 않고 그대로 진행 →", fill=L_GREEN, border=GREEN, size=15, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 8.65, 2.3, .45, .4)
box(s, 9.2, 2.0, 3.4, 1.0, "⑤ 주제 묶기로", fill=WHITE, border=GREEN, size=15, bold=True, align=PP_ALIGN.CENTER)
down_arrow(s, 2.0, 3.15)
box(s, .6, 3.55, 3.2, 1.05, "이미지 일꾼 (옆길)\n한가할 때 내려받아 봄", fill=L_ORANGE, border=ORANGE, size=14, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 3.9, 3.85, .45, .4)
box(s, 4.45, 3.55, 4.1, 1.05, "관련성 판단 (AI의 눈)\n\"이 사진이 글 내용을 보여주나?\"", fill=WHITE, border=PURPLE, size=14, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 8.65, 3.85, .45, .4)
box(s, 9.2, 3.55, 3.4, 1.05, "관련 사진만 한 문장 묘사\n→ 글 좌표에 반영", fill=L_BLUE, border=BLUE, size=13, bold=True, align=PP_ALIGN.CENTER)
box(s, .7, 5.1, 11.9, 1.5,
    "왜 관련성을 따지나요? — 디씨 특성상 첨부의 약 40%는 글과 무관한 반응짤(캐릭터 일러·밈)입니다.\n"
    "무관한 사진 묘사가 좌표에 섞이면 짧은 글일수록 엉뚱한 주제로 끌려갑니다. 실측 예: 전투 기록 스크린샷·순위표는 반영 ✓ / \"내일 아침 추천좀\"에 붙은 벚꽃 일러는 제외 ✗",
    fill=L_GRAY, border=LINE, size=14, bold=True, align=PP_ALIGN.CENTER)

# ── 11. ⑥ 주제 묶기 ────────────────────────────────────────────────
s = base("⑤ 주제 묶기 — '지금 무슨 이야기가 도는가'", "주제는 카페의 원탁입니다. 매분, 최근 24시간의 모든 글을 백지에서 다시 앉힙니다.", "5단계 · 주제")
box(s, .7, 2.0, 5.9, 3.8,
    "합석 규칙 (average-link)\n\n새 글은 원탁에 앉은 '모든 사람과 평균적으로'\n충분히 통해야 합석할 수 있습니다.\n\n"
    "· 6명 미만 원탁은 해산 (잡담으로 간주)\n· 어디에도 못 앉으면 혼자 새 원탁 시작",
    fill=L_GREEN, border=GREEN, size=16, bold=True)
box(s, 6.85, 2.0, 5.9, 3.8,
    "왜 '전원과 평균'인가요?\n\n예전 규칙(원탁의 '분위기 평균'과만 비교)은\n잡탕 원탁일수록 분위기가 두루뭉술해져서\n"
    "서로 안 통하는 156명이 한 원탁에 앉는\n'눈덩이' 사고가 났습니다.\n\n전원-평균 규칙은 잡탕이 될수록 스스로\n합석 문턱이 높아져 눈덩이가 안 생깁니다.",
    fill=WHITE, border=RED, size=15, bold=True)
text(s, .8, 6.1, 11.8, .5, "결과: 최근 창에서 20~30개의 조밀한 원탁 — \"콜라보 카페 방문\", \"루디 반격 불만\", \"닉네임 의무화 논란\" 같은 구체적 사건들.",
     15, NAVY, True, PP_ALIGN.CENTER)

# ── 12. 어제 그 원탁 알아보기 (1:1 매칭) ────────────────────────────
s = base("주제에 이름표 달기 — '어제 그 원탁'을 알아봅니다", "원탁 번호는 매분 새로 매겨지므로, 저장된 주제와 대조해 같은 사건이면 번호를 이어받습니다.", "주제 추적")
box(s, .6, 2.1, 3.5, 2.7, "이번 분(分)의 원탁들\n\nA. 루디 불만 12건\nB. 콜라보 카페 9건\nC. 새 이야기 7건",
    fill=WHITE, border=GREEN, size=15, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 4.2, 3.2, .5, .45)
box(s, 4.8, 2.1, 4.2, 2.7, "1:1 대조 (닮은 순서대로)\n\nA ↔ 주제 #241 (닮음 0.95) ✓\nB ↔ 주제 #151 (닮음 0.97) ✓\nC ↔ 닮은 주제 없음",
    fill=L_BLUE, border=BLUE, size=15, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 9.1, 3.2, .5, .45)
box(s, 9.7, 2.1, 3.0, 2.7, "결과\n\nA·B: 기존 이름·이력 유지\nC: AI가 새 이름 지음",
    fill=L_PURPLE, border=PURPLE, size=15, bold=True, align=PP_ALIGN.CENTER)
box(s, .7, 5.25, 11.9, 1.3,
    "왜 1:1인가요? — 예전에는 여러 원탁이 같은 이름표를 동시에 가져갈 수 있어서, 애써 나눈 원탁들이 도로 합쳐지는 문제가 있었습니다.\n"
    "지금은 가장 닮은 원탁 하나만 이름표를 이어받아, 서로 다른 사건이 한 주제로 섞이지 않습니다.",
    fill=L_GRAY, border=LINE, size=14, bold=True, align=PP_ALIGN.CENTER)

# ── 13. 주제 품질 관리 ─────────────────────────────────────────────
s = base("원탁의 질서 유지 — 세 가지 자동 교정", "벡터 묶기의 한계(비슷한 단어, 다른 얘기)를 AI가 계속 바로잡습니다.", "품질 관리")
cards = [("멤버 검증  (잘못 앉은 사람 내보내기)",
          "AI가 원탁 멤버를 읽고 \"정말 이 주제 얘기인가?\" O/X 판정.\nX 판정 글은 화면·요약·알림에서 제외. 판정은 저장되어 두 번 묻지 않음.", GREEN),
         ("병합  (갈라진 원탁 붙이기)",
          "두 원탁이 아주 많이 닮으면(0.92↑) 자동 병합, 애매하면(0.75~0.92)\nAI가 \"같은 사건인가\" 판정. '별개' 판정은 기억해서 재질문 없음.", BLUE),
         ("재명명  (이름이 내용을 따라가게)",
          "대화가 흘러가서 이름과 내용이 멀어지면 현재 멤버 기준으로 이름을\n다시 지음. 예: 「라드 필수성」 → 실제로는 루디 각성 얘기가 된 경우.", PURPLE)]
for i, (a, b, c) in enumerate(cards):
    y = 2.0 + i * 1.5
    box(s, .7, y, 4.1, 1.25, a, fill=c, color=WHITE, size=15, bold=True, align=PP_ALIGN.CENTER)
    box(s, 4.95, y, 7.65, 1.25, b, fill=WHITE, border=c, size=14, bold=True)
text(s, .8, 6.6, 11.8, .45, "AI 판정(검증·병합)은 전부 파일로 기억해 같은 질문에 두 번 돈을 쓰지 않습니다.", 14, MUTED, False, PP_ALIGN.CENTER)

# ── 14. ⑥ 급증 감지 ────────────────────────────────────────────────
s = base("⑥ 급증 감지 — '많다'가 아니라 '갑자기 뜨겁다'", "주제마다 열기 점수를 매기고, 평소 대비 급증했을 때만 다음 단계로 보냅니다.", "6단계 · 버스트")
box(s, .7, 2.0, 3.75, 3.7,
    "열기(heat) 점수\n\n방금 글 = 1점\n30분 전 글 = 0.5점\n1시간 전 글 = 0.25점\n\n식으면 점수도 스르르 감소\n(반감기 30분)",
    fill=L_ORANGE, border=ORANGE, size=16, bold=True, align=PP_ALIGN.CENTER)
box(s, 4.8, 2.0, 3.75, 3.7,
    "급증(버스트) 판정\n\n두 조건 동시 충족:\n\n① 최근 1시간 5건 이상 (절대량)\n② 그 주제의 평소 시간당\n    평균의 3배 이상 (비율)",
    fill=L_RED, border=RED, size=16, bold=True, align=PP_ALIGN.CENTER)
box(s, 8.9, 2.0, 3.75, 3.7,
    "왜 이렇게 하나요?\n\n큰 주제는 원래 글이 많으니\n'비율'로, 작은 주제의 우연은\n'절대량'으로 거릅니다.\n\n급증 아닌 주제는 비싼\n요약·판단을 아예 안 돌립니다.",
    fill=WHITE, border=LINE, size=15, bold=True, align=PP_ALIGN.CENTER)

# ── 15. ⑦⑧ 판단과 알림 ────────────────────────────────────────────
s = base("⑦ 요약·판단 → ⑧ 알림 — 그리고 '소강'까지", "급증한 주제만 이 파이프에 들어옵니다. 알림 후 이슈가 식으면 종료 신호도 자동으로 나갑니다.", "7·8단계 · 알림")
flow = [("재검증", "이물질 글 제거", GREEN), ("요약", "대표 글 최대 15건\n→ 3~5문장", CYAN),
        ("라벨", "글 읽기 AI 결과\n다수결", PURPLE), ("판단관 O/X", "알릴 사건인가\n한 글자 판정", RED),
        ("알림 발송", "쿨다운 60분\n(도배 방지)", BLUE)]
for i, (a, b, c) in enumerate(flow):
    x = .55 + i * 2.55
    box(s, x, 2.2, 2.2, 1.6, f"{a}\n{b}", fill=WHITE, border=c, size=14, bold=True, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + 2.22, 2.78, .3, .35)
box(s, .7, 4.35, 5.9, 2.0,
    "판단관에게 들어가는 한 줄\n\n\"주제라벨___대표제목___요약\"\n\n학습 때 본 것과 똑같은 형식이라 판단이 안정적.\nO/X 모두 이력에 남고, 담당자의 👍👎 피드백은\n다음 학습 데이터로 쌓입니다.",
    fill=L_PURPLE, border=PURPLE, size=14, bold=True)
box(s, 6.75, 4.35, 5.9, 2.0,
    "소강(이슈 종료) 알림\n\n알림 나간 주제를 계속 감시하다가\n열기가 정점의 25% 아래 + 최근 1시간 5건 미만이\n"
    "10분 연속이면 \"🌤 소강\" 한 번 발송.\n담당자가 상황 종료를 알 수 있습니다.",
    fill=L_GREEN, border=GREEN, size=14, bold=True)

# ── 16. 하루 시나리오 ──────────────────────────────────────────────
s = base("실제로 이렇게 흘러갑니다 — '루디 패치의 날' (실사례)", "2026-07-11 실제 알림 이력을 재구성한 타임라인입니다.", "시나리오")
timeline = [("06~09시", "패치 직후 불만 글이\n띄엄띄엄 유입", "원탁이 조용히 생김\n(주제: 루디 각성 불만)", MUTED),
            ("13시 40분", "1시간에 글 10여 건 몰림\n평소의 3배 초과", "🔥 버스트 감지\n→ 재검증·요약 시작", ORANGE),
            ("13시 59분", "판단관: \"반격 렉·밸런스\n실불만 다수\" → O", "🚨 알림 발송\n요약과 함께 이력 기록", RED),
            ("14시 55분", "유입 급감, 열기가\n정점의 17%까지 하락", "🌤 소강 알림\n관제 사이클 종료", GREEN)]
for i, (t, a, b, c) in enumerate(timeline):
    x = .55 + i * 3.15
    box(s, x, 2.1, 2.85, .6, t, fill=c, color=WHITE, size=16, bold=True, align=PP_ALIGN.CENTER)
    box(s, x, 2.85, 2.85, 1.45, a, fill=WHITE, border=c, size=13, bold=True, align=PP_ALIGN.CENTER)
    box(s, x, 4.45, 2.85, 1.45, b, fill=L_GRAY, border=c, size=13, bold=True, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, x + 2.86, 3.4, .28, .35)
text(s, .8, 6.2, 11.8, .5, "운영자는 알림 두 줄(발생·소강)만 보고도 사건의 시작과 끝을 알 수 있습니다. 상세 화면에서는 당시 글들을 그대로 볼 수 있습니다.",
     14, NAVY, True, PP_ALIGN.CENTER)

# ── 17. 화면 안내 ──────────────────────────────────────────────────
s = base("화면은 이렇게 읽습니다", "브라우저 화면은 10~15초마다 자동 갱신됩니다. 좌측에서 게임을 전환할 수 있습니다.", "화면 안내")
screens = [("실시간 관제 (홈)", "카테고리 열기 보드 · AI 슬롯 상태 · 최근 수집 글", BLUE),
           ("주제", "지금 도는 사건 목록 — 알림 → 확산 → 열기 순 정렬", GREEN),
           ("주제 상세", "기본은 '현재 창'의 글만 표시 · '전체 이력'은 토글로", CYAN),
           ("알림", "발송 판정과 이력 · 담당자 👍👎 피드백(학습 데이터)", RED),
           ("타임머신", "과거 어느 시점이든 당시 대시보드 그대로 재생", PURPLE),
           ("점검 도구", "재분류 · 안전망 점검 · 판단 모델 단독 테스트", ORANGE)]
for i, (a, b, c) in enumerate(screens):
    x = .65 + (i % 2) * 6.25; y = 2.0 + (i // 2) * 1.5
    box(s, x, y, 5.9, 1.25, f"{a}\n{b}", fill=WHITE, border=c, size=15, bold=True)
text(s, .8, 6.6, 11.8, .45, "타임머신: 매분의 화면 상태를 30일간 저장 — \"어제 밤 무슨 일이 있었나\"를 그대로 되감아 봅니다.",
     14, MUTED, False, PP_ALIGN.CENTER)

# ── 18. 조정 다이얼 ────────────────────────────────────────────────
s = base("운영자가 만질 수 있는 다이얼 (현재 값)", "숫자는 이 게시판 실데이터로 보정한 값입니다 — 임베딩 모델이나 게임이 바뀌면 재보정이 필요합니다.", "설정")
rows = [("주제 합석 문턱 0.60", "원탁 전원과의 평균 닮음", "낮추면 잡탕 원탁 / 높이면 원탁이 안 생김"),
        ("주제 인정 최소 6건", "원탁 성립 인원", "낮추면 자잘한 주제 급증 / 높이면 작은 사건 놓침"),
        ("주제 이름표 승계 0.85", "'어제 그 원탁' 판정", "낮추면 다른 사건 흡수 / 높이면 새 주제 남발"),
        ("자동 병합 0.92", "확실한 쌍둥이 원탁 합치기", "낮추면 잘못 합쳐질 위험 (되돌리기 어려움)"),
        ("열기 반감기 30분", "글 한 건의 영향 지속", "짧으면 민감 / 길면 오래 뜨겁게 보임"),
        ("급증 = 5건 & 3배", "알림 후보 문턱", "낮추면 오탐 증가 / 높이면 사건을 늦게 봄"),
        ("재알림 쿨다운 60분", "같은 주제 반복 알림 방지", "짧으면 시끄러움 / 길면 재점화를 놓침")]
for i, (a, b, c) in enumerate(rows):
    y = 1.9 + i * .69
    box(s, .6, y, 3.6, .55, a, fill=L_BLUE, border=BLUE, size=13, bold=True)
    box(s, 4.3, y, 3.4, .55, b, fill=WHITE, border=LINE, size=13, bold=True)
    box(s, 7.8, y, 4.85, .55, c, fill=WHITE, border=LINE, size=12)

# ── 19. 한계 ───────────────────────────────────────────────────────
s = base("결과를 해석할 때 알아둘 것", "이 시스템은 '자동 정답기'가 아니라, 운영자가 빨리 볼 후보를 만들어 주는 관제 도구입니다.", "주의사항")
notes = [("AI는 틀릴 수 있습니다", "분류·이름·O/X 판단 모두 오판 가능. 알림 화면의 👍👎 피드백이 다음 학습의 재료가 됩니다.", RED),
         ("6건 미만은 주제가 안 됩니다", "의미 있는 사건이라도 글이 적으면 원탁이 서지 않습니다. 카테고리 열기 그래프가 보조 신호입니다.", ORANGE),
         ("주제와 카테고리는 1:1이 아닙니다", "한 사건(주제)에 밸런스·콘텐츠 글이 섞일 수 있습니다 — 정상입니다.", BLUE),
         ("숫자는 이 게시판 전용입니다", "합석 문턱·병합 기준 등은 실데이터로 보정한 값 — 다른 게임·모델에 그대로 옮기면 안 맞습니다.", PURPLE),
         ("이미지 판정도 완벽하지 않습니다", "관련/무관 경계 사례(캐릭터 일러 등)는 가끔 틀립니다. 임베딩에만 쓰여 영향 범위는 제한적입니다.", GREEN)]
for i, (a, b, c) in enumerate(notes):
    y = 1.95 + i * .95
    box(s, .7, y, 3.9, .8, a, fill=WHITE, border=c, size=14, bold=True, align=PP_ALIGN.CENTER)
    box(s, 4.75, y, 7.85, .8, b, fill=WHITE, border=LINE, size=13)

# ── 20. 마무리 ─────────────────────────────────────────────────────
s = base("기억할 세 문장", section="요약")
box(s, .9, 2.0, 11.55, 1.25, "1. 모든 글은 AI가 한 번씩 읽어 7개 서랍(카테고리)에 넣습니다 — 게시판의 지형.",
    fill=L_BLUE, border=BLUE, size=20, bold=True, align=PP_ALIGN.CENTER)
box(s, .9, 3.5, 11.55, 1.25, "2. 같은 사건의 글은 원탁(주제)으로 묶고, 검증·병합·재명명으로 질서를 유지합니다 — 사건의 단위.",
    fill=L_GREEN, border=GREEN, size=19, bold=True, align=PP_ALIGN.CENTER)
box(s, .9, 5.0, 11.55, 1.25, "3. 갑자기 뜨거워진 원탁만 요약해 O/X 판단을 거쳐 알리고, 식으면 소강까지 알립니다 — 관제의 완결.",
    fill=L_PURPLE, border=PURPLE, size=19, bold=True, align=PP_ALIGN.CENTER)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"저장: {OUT} ({len(prs.slides)}장)")
